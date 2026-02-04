from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Paleta de colores: Teal Trust (profesional y confiable)
PRIMARY_COLOR = RGBColor(2, 128, 144)      # Teal #028090
SECONDARY_COLOR = RGBColor(0, 168, 150)    # Seafoam #00A896
ACCENT_COLOR = RGBColor(2, 195, 154)       # Mint #02C39A
DARK_BG = RGBColor(30, 39, 97)             # Navy oscuro #1E2761
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(242, 242, 242)
TEXT_DARK = RGBColor(54, 54, 54)
TEXT_LIGHT = RGBColor(100, 116, 139)

def add_title_slide(prs, title, subtitle=""):
    """Slide de título con diseño atractivo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Fondo oscuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Línea decorativa
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3), Inches(3), Inches(4), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    return slide

def add_section_slide(prs, section_title):
    """Slide divisor de sección"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo con gradiente simulado (dos rectángulos)
    bg1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY_COLOR
    bg1.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, numbered=False):
    """Slide de contenido con bullets o texto"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo claro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Barra superior decorativa
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_box.text_frame.margin_top = Inches(0)
    title_box.text_frame.margin_bottom = Inches(0)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(4))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        if isinstance(item, tuple):  # (header, text)
            header, text = item
            p.text = f"{header}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_after = Pt(6)
            
            # Agregar texto descriptivo
            p2 = text_frame.add_paragraph()
            p2.text = text
            p2.font.size = Pt(14)
            p2.font.color.rgb = TEXT_DARK
            p2.space_after = Pt(12)
            p2.level = 1
        else:
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)
            if numbered:
                p.level = 0
            else:
                p.level = 0
    
    return slide

def add_quote_slide(prs, quote_text):
    """Slide de cita destacada"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Barra lateral izquierda
    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = ACCENT_COLOR
    left_bar.line.fill.background()
    
    # Caja de cita
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
    text_frame = quote_box.text_frame
    text_frame.text = f'"{quote_text}"'
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============= CREAR SLIDES =============

# 1. Portada
add_title_slide(prs, "KICK OFF 2026", "Operaciones Data GIS - Servinformación")

# 2. Sección: Reflexión y Aprendizaje
add_section_slide(prs, "Reflexión y Aprendizaje 2025")

# 3. Lecciones aprendidas 2025
add_content_slide(prs, "¿Qué lecciones aprendidas deja el 2025?", [
    "2025 marcó un cambio irreversible: la convergencia de IA, RPA y soluciones a medida redefinió nuestra capacidad operativa.",
    "",
    "La lección fundamental: entender la automatización como una estrategia de Inteligencia Aumentada, no solo como programación.",
    "",
    "La IA transforma operaciones de 'procesos manuales' a 'validación asistida', potenciando el talento humano.",
    "",
    "Aprendizaje crítico: La IA no corrige el caos, lo amplifica. Los modelos de ML exigieron estandarización más rigurosa.",
    "",
    "El éxito radica en preparar nuestra data y cultura operativa para trabajar en simbiosis con estas tecnologías."
])

# 4. Estrategias que marcaron diferencia
add_content_slide(prs, "Estrategias que marcaron la diferencia en 2025", [
    ("1. Células de Desarrollo Híbrido", 
     "Integramos capacidades de programación (Python/SQL) en equipos de operación. Empoderamos líderes técnicos GIS para crear soluciones ágiles in-house, reduciendo drásticamente tiempos de respuesta."),
    ("2. Modularización de Flujos de Trabajo",
     "Dividimos procesos masivos en micro-tareas estandarizadas. Aplicamos automatización quirúrgica donde aporta valor, enfocando talento humano en QA/QC y análisis complejo.")
])

# 5. Indicadores de gestión
add_content_slide(prs, "Indicadores de Gestión 2024-2025", [
    ("1. Productividad Efectiva por Analista",
     "Unidades GIS procesadas/aprobadas por hora. Evidencia el aumento al sustituir digitación manual por validación asistida."),
    ("2. Calidad en Primera Entrega (FTY)",
     "% de lotes sin errores en primer intento. Refleja cómo la automatización 'blindó' la operación."),
    ("3. Tasa de Retrabajo Operativo",
     "% de producción devuelto para correcciones. Cuantifica la reducción del desperdicio operativo.")
])

# 6. Sección: Estrategia y Visión 2026
add_section_slide(prs, "Estrategia y Visión 2026")

# 7. Metodologías para optimizar
add_content_slide(prs, "Metodologías para mayor agilidad", [
    ("Gestión por Casos de Uso",
     "Organizar trabajo desde la decisión de negocio del cliente, no desde la tecnología. Elimina análisis innecesarios."),
    ("Arquitectura Modular GIS",
     "Componentes reutilizables que se ensamblan según necesidad. Reduce tiempos de implementación drásticamente."),
    ("Metodologías Ágiles Adaptadas",
     "Scrum (sprints 2 semanas) • Kanban (BPO/Censos) • DevOps GIS (integración continua)")
])

# 8. Cambios estructurales 2026
add_content_slide(prs, "Cambios estructurales para 2026", [
    ("1. Centro de Excelencia en Automatización",
     "Estandarizar y escalar scripts automatizados hacia las 4 subáreas, eliminando duplicidad."),
    ("2. Equipos Multifuncionales Cross-Área",
     "Combinar expertise de Cartografía Latam, Colombia, BPO y Censos para proyectos complejos."),
    ("3. Arquitectura de Datos Modular",
     "Componentes reutilizables y pre-validados que aceleran nuevos proyectos sin sacrificar calidad.")
])

# 9. Pilar estratégico: Data confiable
add_content_slide(prs, "Pilar Estratégico: Data como Producto Vivo", [
    ("Productos de Data Vivos",
     "Actualización continua en horas, no por lotes. Versiones claras que evitan reprocesos."),
    ("Calidad by Design",
     "Validaciones automatizadas en cada etapa. Métricas de precisión, consistencia y actualidad."),
    ("Trazabilidad Total",
     "Documentación automática de origen, transformaciones y validaciones. Confianza verificable."),
    ("Diseño para Consumo Inmediato",
     "Data estructurada para analítica e integración directa. Ventanas de frescura medibles.")
])

# 10. Sección: Innovación y Adaptación
add_section_slide(prs, "Innovación y Adaptación al Mercado")

# 11. Soluciones Ready-to-use
add_content_slide(prs, "Tecnologías para soluciones 'Ready-to-use'", [
    ("IA para Interpretación Geoespacial",
     "Guiar lectura de mapas, explicar patrones, sugerir oportunidades. IA como copiloto del territorio."),
    ("Analítica Integrada como Estándar",
     "Cliente recibe data + indicadores + interpretación lista para decisiones desde día 1."),
    ("Automatización Inteligente",
     "Procesos repetitivos automatizados. Entrega rápida, consistente y escalable en 48-72 horas.")
])

# 12. Servicios consultivos de valor
add_content_slide(prs, "Servicios Consultivos de Alto Valor", [
    ("Interpretación Territorial Asistida",
     "Monitoreo continuo de variables críticas. Alertas y recomendaciones estratégicas proactivas."),
    ("Análisis por Escenarios Parametrizables",
     "Co-creación de soluciones únicas embebidas con clientes. Scoring territorial predictivo personalizado."),
    ("Validación para Decisiones Críticas",
     "Certificación de data geoespacial. Trazabilidad, calidad demostrable, cumplimiento regulatorio."),
    ("Industry Vertical Solutions",
     "Expertise en Telcos, Retail, Gobierno, Utilities. ROI medible y verificable.")
])

# 13. Compromiso innegociable (slide final destacado)
add_quote_slide(prs, "Convertir la información geoespacial en decisiones claras, confiables y accionables, entregadas con la velocidad que exige el mercado y el nivel de calidad que respalda el crecimiento de Servinformación.")

# 14. Slide de cierre
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

# Texto final
final_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
text_frame = final_box.text_frame
text_frame.text = "2026: El Año del Crecimiento"
p = text_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
text_frame2 = subtitle_box.text_frame
text_frame2.text = "Operaciones Data GIS • Servinformación"
p2 = text_frame2.paragraphs[0]
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT_COLOR
p2.alignment = PP_ALIGN.CENTER

# Guardar presentación
prs.save('/home/claude/Kick_Off_2026_Presentacion.pptx')
print("✓ Presentación creada exitosamente")
